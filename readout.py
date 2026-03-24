import streamlit as st
import pandas as pd
import json
import re
import glob
from pathlib import Path
from concurrent.futures import ThreadPoolExecutor, as_completed
from pydantic import BaseModel, Field
from typing import Optional
import litellm
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import io
import tempfile
import math
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import FancyBboxPatch

# ── Pydantic model ──────────────────────────────────────────────────────────────

class CompetencyResult(BaseModel):
    competency: str
    score: int = Field(..., ge=1, le=5)
    rationale: str

# ── Helpers ─────────────────────────────────────────────────────────────────────

SCORE_EMOJI = {1: "🔴", 2: "🟠", 3: "🟡", 4: "🟢", 5: "✅"}

def load_prompt_files(prompt_dir: str) -> dict[str, str]:
    """Load all .md files from the prompt directory."""
    prompts = {}
    for path in sorted(Path(prompt_dir).glob("*.md")):
        prompts[path.stem] = path.read_text(encoding="utf-8")
    return prompts

def load_csv_robust(file) -> pd.DataFrame:
    """Try increasingly lenient CSV parsing strategies until one works."""
    raw = file.read()
    file.seek(0)

    strategies = [
        dict(encoding="utf-8"),
        dict(encoding="utf-8-sig"),          # BOM
        dict(encoding="latin-1"),
        dict(encoding="utf-8",  sep=None, engine="python"),   # auto-detect separator
        dict(encoding="latin-1", sep=None, engine="python"),
        dict(encoding="utf-8",  on_bad_lines="skip"),         # skip malformed rows
        dict(encoding="latin-1", on_bad_lines="skip"),
        dict(encoding="utf-8",  on_bad_lines="skip", sep=None, engine="python"),
    ]

    last_err = None
    for kwargs in strategies:
        try:
            buf = io.BytesIO(raw)
            df = pd.read_csv(buf, **kwargs)
            if df.empty:
                continue
            return df
        except Exception as e:
            last_err = e

    raise ValueError(f"Could not parse CSV after all strategies. Last error: {last_err}")


def csv_to_text(df: pd.DataFrame) -> str:
    return df.to_csv(index=False)

def run_litellm_call(
    competency_name: str,
    prompt_template: str,
    chat_log_text: str,
    model: str,
    api_key: str,
) -> CompetencyResult:
    """Inject chat log into prompt and call LiteLLM; parse result into pydantic model."""
    prompt = prompt_template.replace("{{chat_log}}", chat_log_text)

    system = (
        "You are an expert evaluator. Respond ONLY with a valid JSON object matching "
        "this schema: {\"competency\": string, \"score\": int (1-5), \"rationale\": string}. "
        "No markdown fences, no extra text."
    )

    response = litellm.completion(
        model=model,
        api_key=api_key or None,
        messages=[
            # {"role": "system", "content": system},
            {"role": "user", "content": prompt},
        ],
        temperature=0,
        top_p=0.0,
        response_format=CompetencyResult,
    )

    raw = response.choices[0].message.content.strip()
    raw = re.sub(r"^```(?:json)?|```$", "", raw, flags=re.MULTILINE).strip()
    data = json.loads(raw)
    data["competency"] = competency_name
    return CompetencyResult(**data)


def evaluate_chat_log(
    filename: str,
    df: pd.DataFrame,
    prompts: dict[str, str],
    model: str,
    api_key: str,
    max_workers: int,
) -> dict[str, CompetencyResult | str]:
    """Run all competency prompts in parallel threads for one CSV."""
    chat_text = csv_to_text(df)
    results: dict[str, CompetencyResult | str] = {}

    with ThreadPoolExecutor(max_workers=max_workers) as executor:
        futures = {
            executor.submit(
                run_litellm_call, name, tpl, chat_text, model, api_key
            ): name
            for name, tpl in prompts.items()
        }
        for future in as_completed(futures):
            name = futures[future]
            try:
                results[name] = future.result()
            except Exception as exc:
                results[name] = f"ERROR: {exc}"

    return results


def build_radar_chart(
    filename: str,
    comp_results: dict,
    competency_names: list[str],
) -> io.BytesIO:
    """Render a radar chart for one chat log; return PNG bytes in a BytesIO buffer."""
    scores = []
    labels = []
    for comp in competency_names:
        res = comp_results.get(comp)
        scores.append(res.score if isinstance(res, CompetencyResult) else 0)
        labels.append(comp)

    N = len(labels)
    angles = [2 * math.pi * i / N for i in range(N)]
    angles += angles[:1]          # close the loop
    scores_plot = scores + scores[:1]

    fig, ax = plt.subplots(figsize=(5, 5), subplot_kw=dict(polar=True))
    fig.patch.set_facecolor("#F8F9FA")
    ax.set_facecolor("#F8F9FA")

    # Grid rings
    ax.set_ylim(0, 5)
    ax.set_yticks([1, 2, 3, 4, 5])
    ax.set_yticklabels(["1", "2", "3", "4", "5"],
                       fontsize=7, color="#ADB5BD")
    ax.yaxis.set_tick_params(pad=2)

    # Spoke labels
    ax.set_xticks(angles[:-1])
    wrapped = ["\n".join(l.split()) if len(l) > 12 else l for l in labels]
    ax.set_xticklabels(wrapped, fontsize=9, color="#212529", fontweight="bold")

    # Grid styling
    ax.grid(color="#DEE2E6", linestyle="--", linewidth=0.8, alpha=0.9)
    ax.spines["polar"].set_color("#CED4DA")

    # Fill
    ax.plot(angles, scores_plot, color="#4361EE", linewidth=2.2, linestyle="solid")
    ax.fill(angles, scores_plot, color="#4361EE", alpha=0.20)

    # Score dots
    for ang, sc in zip(angles[:-1], scores):
        color = ["#DC2626","#EA580C","#CA8A04","#16A34A","#059669"][sc - 1] if sc else "#ADB5BD"
        ax.plot(ang, sc, "o", color=color, markersize=8, zorder=5)

    average = sum(scores) / len(scores) if scores else 0
    ax.set_title(Path(filename).stem + f" (Average: {average:.2f}/5)", fontsize=12, fontweight="bold",
                 color="#212529", pad=16)
    

    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight",
                facecolor=fig.get_facecolor())
    plt.close(fig)
    buf.seek(0)
    return buf


# ── PowerPoint export ─────────────────────────────────────────────────────────

def build_pptx(all_results: dict[str, dict], competency_names: list[str],
               radar_buffers: dict[str, io.BytesIO]) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # ── Palette ──────────────────────────────────────────────────────────────────
    WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
    OFF_WHITE   = RGBColor(0xF8, 0xF9, 0xFA)
    LIGHT_GREY  = RGBColor(0xE9, 0xEC, 0xEF)
    MID_GREY    = RGBColor(0xAD, 0xB5, 0xBD)
    DARK_TEXT   = RGBColor(0x21, 0x25, 0x29)
    SUBTLE_TEXT = RGBColor(0x49, 0x54, 0x57)
    ACCENT      = RGBColor(0x43, 0x61, 0xEE)   # vivid indigo — title bar & header
    ACCENT_SOFT = RGBColor(0xEB, 0xEF, 0xFF)   # tint for alternating rows
    SCORE_COLORS = {
        1: RGBColor(0xDC, 0x26, 0x26),   # red-600
        2: RGBColor(0xEA, 0x58, 0x0C),   # orange-600
        3: RGBColor(0xCA, 0x8A, 0x04),   # yellow-600
        4: RGBColor(0x16, 0xA3, 0x4A),   # green-600
        5: RGBColor(0x05, 0x96, 0x69),   # emerald-600
    }
    SCORE_BG = {
        1: RGBColor(0xFE, 0xE2, 0xE2),
        2: RGBColor(0xFF, 0xED, 0xD5),
        3: RGBColor(0xFF, 0xF9, 0xC2),
        4: RGBColor(0xDC, 0xFC, 0xE7),
        5: RGBColor(0xD1, 0xFA, 0xEE),
    }

    # ── Helpers ──────────────────────────────────────────────────────────────────
    def rect(slide, l, t, w, h, fill, line_color=None, radius=False):
        shp = slide.shapes.add_shape(1, l, t, w, h)
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
        if line_color:
            shp.line.color.rgb = line_color
            shp.line.width = Pt(0.5)
        else:
            shp.line.fill.background()
        return shp

    def label(slide, text, l, t, w, h,
              size=11, bold=False, color=DARK_TEXT,
              align=PP_ALIGN.LEFT, italic=False):
        txb = slide.shapes.add_textbox(l, t, w, h)
        txb.word_wrap = True
        tf = txb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        run.font.name = "Calibri"
        return txb

    # ── Slide 1: Summary table ───────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    rect(slide, 0, 0, prs.slide_width, prs.slide_height, OFF_WHITE)

    # Title bar
    rect(slide, 0, 0, prs.slide_width, Inches(0.85), ACCENT)
    label(slide, "Chat Log Competency Evaluation",
          Inches(0.35), Inches(0.15), prs.slide_width - Inches(0.7), Inches(0.6),
          size=22, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    filenames  = list(all_results.keys())
    n_files    = len(filenames)
    label_col  = Inches(2.6)
    score_col  = Inches(1.6)
    row_h      = Inches(0.46)
    start_x    = Inches(0.35)
    start_y    = Inches(1.05)
    table_w    = label_col + score_col * n_files

    # Column headers
    rect(slide, start_x, start_y, label_col, row_h, ACCENT)
    label(slide, "COMPETENCY", start_x + Inches(0.12), start_y + Inches(0.1),
          label_col - Inches(0.15), row_h, size=9, bold=True, color=WHITE)

    for j, fn in enumerate(filenames):
        cx = start_x + label_col + j * score_col
        rect(slide, cx, start_y, score_col, row_h, ACCENT,
             line_color=RGBColor(0x6B, 0x82, 0xF4))
        label(slide, Path(fn).stem[:20], cx + Inches(0.06), start_y + Inches(0.08),
              score_col - Inches(0.1), row_h, size=9, bold=True,
              color=WHITE, align=PP_ALIGN.CENTER)

    # Data rows
    for i, comp in enumerate(competency_names):
        ry = start_y + (i + 1) * row_h
        row_fill = WHITE if i % 2 == 0 else ACCENT_SOFT
        rect(slide, start_x, ry, label_col, row_h, row_fill,
             line_color=LIGHT_GREY)
        label(slide, comp, start_x + Inches(0.12), ry + Inches(0.1),
              label_col - Inches(0.15), row_h, size=10, color=DARK_TEXT)

        for j, fn in enumerate(filenames):
            cx = start_x + label_col + j * score_col
            rect(slide, cx, ry, score_col, row_h, row_fill, line_color=LIGHT_GREY)
            res = all_results[fn].get(comp)
            if isinstance(res, CompetencyResult):
                bg  = SCORE_BG.get(res.score, LIGHT_GREY)
                fg  = SCORE_COLORS.get(res.score, DARK_TEXT)
                pill_w, pill_h = Inches(0.55), Inches(0.28)
                pill_x = cx + (score_col - pill_w) / 2
                pill_y = ry + (row_h - pill_h) / 2
                rect(slide, pill_x, pill_y, pill_w, pill_h, bg, line_color=fg)
                label(slide, f"{res.score} / 5", pill_x, pill_y,
                      pill_w, pill_h, size=9, bold=True, color=fg,
                      align=PP_ALIGN.CENTER)
            else:
                label(slide, "ERR", cx, ry, score_col, row_h,
                      size=9, color=SCORE_COLORS[1], align=PP_ALIGN.CENTER)

    # Bottom rule
    rect(slide, start_x, start_y + (len(competency_names) + 1) * row_h,
         table_w, Pt(2), ACCENT)

    # ── One slide per file — rationale cards ─────────────────────────────────────
    for fn, comp_results in all_results.items():
        slide = prs.slides.add_slide(blank_layout)
        rect(slide, 0, 0, prs.slide_width, prs.slide_height, OFF_WHITE)
        rect(slide, 0, 0, prs.slide_width, Inches(0.85), ACCENT)
        label(slide, f"Rationales  ·  {Path(fn).stem}",
              Inches(0.35), Inches(0.15), prs.slide_width - Inches(0.7), Inches(0.6),
              size=20, bold=True, color=WHITE)

        cols   = 2
        card_w = Inches(6.2)
        card_h = Inches(1.55)
        gap    = Inches(0.18)
        pad_l  = (prs.slide_width - cols * card_w - (cols - 1) * gap) / 2

        for idx, comp in enumerate(competency_names):
            col_i, row_i = idx % cols, idx // cols
            cx = pad_l + col_i * (card_w + gap)
            cy = Inches(1.05) + row_i * (card_h + gap)
            if cy + card_h > prs.slide_height - Inches(0.15):
                break

            res = comp_results.get(comp)
            score = res.score if isinstance(res, CompetencyResult) else None
            bar_color = SCORE_COLORS.get(score, MID_GREY)
            bg_color  = SCORE_BG.get(score, LIGHT_GREY) if score else LIGHT_GREY

            # Card base
            rect(slide, cx, cy, card_w, card_h, WHITE, line_color=LIGHT_GREY)
            # Left accent bar
            rect(slide, cx, cy, Inches(0.08), card_h, bar_color)
            # Subtle header tint
            rect(slide, cx + Inches(0.08), cy, card_w - Inches(0.08), Inches(0.4), bg_color)

            if isinstance(res, CompetencyResult):
                # Competency name + score badge on same line
                label(slide, comp.upper(),
                      cx + Inches(0.18), cy + Inches(0.07),
                      card_w - Inches(0.9), Inches(0.3),
                      size=9, bold=True, color=DARK_TEXT)
                # Score badge (top-right of card)
                bdg_w, bdg_h = Inches(0.52), Inches(0.26)
                bdg_x = cx + card_w - bdg_w - Inches(0.12)
                bdg_y = cy + Inches(0.07)
                rect(slide, bdg_x, bdg_y, bdg_w, bdg_h,
                     bar_color, line_color=bar_color)
                label(slide, f"{res.score}/5", bdg_x, bdg_y, bdg_w, bdg_h,
                      size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
                # Rationale body
                label(slide, res.rationale,
                      cx + Inches(0.18), cy + Inches(0.44),
                      card_w - Inches(0.28), card_h - Inches(0.52),
                      size=9, color=SUBTLE_TEXT)
            else:
                label(slide, f"{comp}: {res}",
                      cx + Inches(0.18), cy + Inches(0.12),
                      card_w - Inches(0.28), card_h - Inches(0.2),
                      size=9, color=SCORE_COLORS[1])

    # ── Radar chart slides (max 3 per slide, evenly distributed) ─────────────────
    MAX_PER_SLIDE = 3
    items = list(radar_buffers.items())
    n = len(items)

    # Decide how many slides and how many charts each gets, spread as evenly as possible.
    # e.g. 4 charts → 2 slides of 2; 5 → slides of 3+2; 6 → 2 slides of 3
    n_slides = math.ceil(n / MAX_PER_SLIDE)
    base, extra = divmod(n, n_slides)          # base per slide; first `extra` slides get +1
    distribution = [base + (1 if i < extra else 0) for i in range(n_slides)]

    chart_w = Inches(4.0)
    chart_h = Inches(4.0)
    chart_area_top = Inches(1.0)
    chart_area_h = prs.slide_height - chart_area_top - Inches(0.2)

    item_idx = 0
    for slide_i, per_slide in enumerate(distribution):
        radar_slide = prs.slides.add_slide(blank_layout)
        rect(radar_slide, 0, 0, prs.slide_width, prs.slide_height, OFF_WHITE)
        rect(radar_slide, 0, 0, prs.slide_width, Inches(0.85), ACCENT)
        label(radar_slide,
              f"Competency Radar Charts ({slide_i + 1}/{n_slides})",
              Inches(0.35), Inches(0.15), prs.slide_width - Inches(0.7), Inches(0.6),
              size=22, bold=True, color=WHITE)

        # Evenly space `per_slide` charts horizontally
        total_chart_w = per_slide * chart_w
        h_gap = (prs.slide_width - total_chart_w) / (per_slide + 1)
        cy = chart_area_top + (chart_area_h - chart_h) / 2   # vertically centred

        for pos in range(per_slide):
            fn, buf = items[item_idx]
            cx = h_gap + pos * (chart_w + h_gap)
            buf.seek(0)
            radar_slide.shapes.add_picture(buf, cx, cy, chart_w, chart_h)
            item_idx += 1

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ── Markdown export ──────────────────────────────────────────────────────────────

def build_markdown(all_results: dict, competency_names: list[str]) -> str:
    filenames = list(all_results.keys())
    header = "| Competency | " + " | ".join(Path(f).stem for f in filenames) + " |"
    sep    = "|---|" + "|".join(["---"] * len(filenames)) + "|"
    rows   = []
    for comp in competency_names:
        cells = []
        for fn in filenames:
            res = all_results[fn].get(comp)
            if isinstance(res, CompetencyResult):
                cells.append(f"{SCORE_EMOJI.get(res.score, '')} {res.score}/5")
            else:
                cells.append("Error")
        rows.append(f"| {comp} | " + " | ".join(cells) + " |")

    md = ["# Competency Evaluation Results\n", header, sep, *rows, ""]

    for fn, comp_results in all_results.items():
        md.append(f"\n## {Path(fn).stem}\n")
        for comp in competency_names:
            res = comp_results.get(comp)
            if isinstance(res, CompetencyResult):
                md.append(f"### {comp} — {res.score}/5\n{res.rationale}\n")
            else:
                md.append(f"### {comp}\n{res}\n")

    return "\n".join(md)


# ── Streamlit UI ─────────────────────────────────────────────────────────────────

st.set_page_config(page_title="Uptime Labs Competency Readout Generator", layout="wide")
st.title("Uptime Labs Competency Readout Generator")
st.caption("Upload CSVs · load prompts from Markdown files · score with LiteLLM · export results")

# ── Sidebar config ───────────────────────────────────────────────────────────────
with st.sidebar:
    st.header("Configuration")
    model = st.text_input("LiteLLM model", value="gpt-4.1",
                          help="e.g. gpt-4.1, anthropic/claude-3-5-haiku-20241022, ollama/llama3")
    api_key = st.text_input("API Key (optional if set in env)", type="password")
    max_workers = st.slider("Parallel threads per file", 1, 10, 4)
    prompt_dir = st.text_input("Prompt directory", value="prompts",
                               help="Folder containing one .md file per competency")
    st.divider()
    st.markdown("**Prompt variable:** use `{{chat_log}}` in your `.md` files — it will be replaced with the CSV contents.")

# ── Prompt file management ───────────────────────────────────────────────────────
with st.expander("Prompt files (Markdown)", expanded=True):
    prompt_path = Path(prompt_dir)
    if not prompt_path.exists():
        st.warning(f"Directory `{prompt_dir}` not found. Creating it with a sample prompt.")
        prompt_path.mkdir(parents=True, exist_ok=True)
        sample = (
            "# Empathy Evaluation\n\n"
            "Review the following chat log and evaluate the agent's **empathy** "
            "toward the customer.\n\n"
            "```\n{{chat_log}}\n```\n\n"
            "Score the empathy on a scale of 1–5 and provide a concise rationale."
        )
        (prompt_path / "empathy.md").write_text(sample)

    prompts = load_prompt_files(prompt_dir)

    if prompts:
        selected = st.selectbox("Preview prompt", list(prompts.keys()))
        st.markdown(prompts[selected])
        st.info(f"Loaded **{len(prompts)}** competency prompt(s): {', '.join(prompts.keys())}")
    else:
        st.warning(f"No `.md` files found in `{prompt_dir}/`.")

    uploaded_prompt = st.file_uploader("Upload a new prompt .md file", type=["md"])
    if uploaded_prompt:
        dest = prompt_path / uploaded_prompt.name
        dest.write_bytes(uploaded_prompt.read())
        st.success(f"Saved `{uploaded_prompt.name}` to `{prompt_dir}/`. Reload to refresh.")

# ── CSV upload ───────────────────────────────────────────────────────────────────
st.divider()
csv_files = st.file_uploader(
    "Upload chat log CSV files", type=["csv"], accept_multiple_files=True
)

if csv_files:
    st.subheader("Loaded Files")
    preview_file = st.selectbox("Preview CSV", [f.name for f in csv_files])
    for f in csv_files:
        if f.name == preview_file:
            try:
                df_preview = load_csv_robust(f)
                f.seek(0)
                st.dataframe(df_preview.head(5), use_container_width=True)
            except ValueError as e:
                st.error(f"Could not preview file: {e}")
            break

# ── Run evaluation ───────────────────────────────────────────────────────────────
st.divider()
run_btn = st.button("Run Evaluation", type="primary",
                    disabled=not csv_files or not prompts)

if run_btn:
    all_results: dict[str, dict] = {}
    radar_buffers: dict[str, io.BytesIO] = {}
    competency_names = list(prompts.keys())

    progress = st.progress(0, text="Starting…")
    status_cols = st.columns(len(csv_files))

    for file_idx, csv_file in enumerate(csv_files):
        csv_file.seek(0)
        try:
            df = load_csv_robust(csv_file)
        except ValueError as parse_err:
            status_cols[file_idx].error(f"❌ {csv_file.name}: {parse_err}")
            continue
        status_cols[file_idx].info(f"⏳ {csv_file.name}")

        with st.spinner(f"Evaluating **{csv_file.name}**…"):
            result = evaluate_chat_log(
                filename=csv_file.name,
                df=df,
                prompts=prompts,
                model=model,
                api_key=api_key,
                max_workers=max_workers,
            )
        all_results[csv_file.name] = result
        radar_buffers[csv_file.name] = build_radar_chart(
            csv_file.name, result, competency_names
        )
        status_cols[file_idx].success(f"✅ {csv_file.name}")
        progress.progress((file_idx + 1) / len(csv_files),
                          text=f"Completed {file_idx + 1}/{len(csv_files)} files")

    progress.empty()
    st.success("🎉 Evaluation complete!")

    # ── Summary table ────────────────────────────────────────────────────────────
    st.subheader("Summary Table")
    table_data = {"Competency": competency_names}
    for fn, comp_results in all_results.items():
        short = Path(fn).stem
        table_data[short] = []
        for comp in competency_names:
            res = comp_results.get(comp)
            if isinstance(res, CompetencyResult):
                table_data[short].append(
                    f"{SCORE_EMOJI.get(res.score, '')} {res.score}/5"
                )
            else:
                table_data[short].append("⚠️ Error")

    summary_df = pd.DataFrame(table_data)
    st.dataframe(summary_df, use_container_width=True, hide_index=True)

    st.subheader("Radar Charts")
    radar_cols = st.columns(min(len(radar_buffers), 3))
    for i, (fn, rbuf) in enumerate(radar_buffers.items()):
        rbuf.seek(0)
        radar_cols[i % 3].image(rbuf, use_container_width=True)

    # ── Rationales expandable ─────────────────────────────────────────────────
    st.subheader("Detailed Rationales")
    for fn, comp_results in all_results.items():
        with st.expander(f"{fn}"):
            for comp in competency_names:
                res = comp_results.get(comp)
                if isinstance(res, CompetencyResult):
                    color = ["red","orange","goldenrod","green","teal"][res.score - 1]
                    st.markdown(
                        f"**{comp}** — "
                        f"<span style='color:{color};font-weight:bold'>{SCORE_EMOJI[res.score]} {res.score}/5</span>",
                        unsafe_allow_html=True,
                    )
                    st.caption(res.rationale)
                else:
                    st.error(f"**{comp}**: {res}")
                st.divider()

    # ── Exports ──────────────────────────────────────────────────────────────────
    st.subheader("⬇Export Results")
    ecol1, ecol2 = st.columns(2)

    with ecol1:
        try:
            pptx_bytes = build_pptx(all_results, competency_names, radar_buffers)
            st.download_button(
                label="Download PowerPoint (.pptx)",
                data=pptx_bytes,
                file_name="competency_evaluation.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
        except Exception as e:
            st.error(f"PowerPoint generation failed: {e}")

    with ecol2:
        md_str = build_markdown(all_results, competency_names)
        st.download_button(
            label="Download Markdown (.md)",
            data=md_str.encode(),
            file_name="competency_evaluation.md",
            mime="text/markdown",
        )